"""Builds LULU_Investment_Pitch_Deck.pptx from scratch, following the GIS
Investment Research pitch template (structure + formatting conventions):
Garamond throughout, 0.3" header rectangles at 15pt, a one-sentence descriptor
per slide (no trailing periods), theme colors only, and cited sources.

Recommendation: LONG / OVERWEIGHT on lululemon athletica (NASDAQ: LULU).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import data as D

OUT = os.path.join(os.path.dirname(__file__), "..", "LULU_Investment_Pitch_Deck.pptx")

# ---- Theme colors (used consistently throughout) ----
NAVY = RGBColor(0x1F, 0x2A, 0x44)
CARD = RGBColor(0x99, 0x00, 0x00)   # USC cardinal
GOLD = RGBColor(0xFF, 0xC7, 0x2C)   # USC gold
INK = RGBColor(0x26, 0x26, 0x26)
GREY = RGBColor(0x8C, 0x8C, 0x8C)
LGREY = RGBColor(0xF0, 0xF1, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x7A, 0x3C)
FONT = "Garamond"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, size, color=INK, bold=False, italic=False, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # ensure east-asian/complex also Garamond
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf


def add_para(tf, text, size=14, color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             bullet=False, level=0, space_after=4, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    r = p.add_run()
    r.text = text
    _set_font(r, size, color, bold, italic)
    if bullet:
        _add_bullet(p, color)
    else:
        _no_bullet(p)
    return p


def _no_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    for tag in ('a:buChar', 'a:buAutoNum'):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    none = pPr.makeelement(qn('a:buNone'), {})
    pPr.append(none)


def _add_bullet(p, color):
    pPr = p._p.get_or_add_pPr()
    pPr.set('indent', str(Pt(-12)))
    pPr.set('marL', str(Pt(14)))
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': FONT})
    pPr.append(buFont)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)


def rect(slide, l, t, w, h, fill=NAVY, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def slide_base(title, descriptor, sources=None, page=None):
    """Standard content slide chrome: running header, 0.3in title rectangle at
    15pt, one-sentence descriptor (no trailing period), footer + source."""
    s = prs.slides.add_slide(BLANK)
    # running header (left)
    tb, tf = textbox(s, Inches(0.45), Inches(0.18), Inches(6.5), Inches(0.55))
    add_para(tf, "Investment Research Division", 11, CARD, bold=True, first=True, space_after=0)
    add_para(tf, "lululemon athletica (NASDAQ: LULU)", 12.5, NAVY, bold=True, space_after=0)
    # header rectangle (height 0.3in), title 15pt
    hr = rect(s, Inches(6.0), Inches(0.22), Inches(6.9), Inches(0.3), fill=NAVY)
    htf = hr.text_frame
    htf.word_wrap = True
    htf.margin_top = Pt(0); htf.margin_bottom = Pt(0); htf.margin_right = Pt(6)
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.RIGHT
    hr_run = hp.add_run(); hr_run.text = title
    _set_font(hr_run, 15, WHITE, bold=True)
    # descriptor sentence (no period)
    tb2, tf2 = textbox(s, Inches(0.45), Inches(0.62), Inches(12.45), Inches(0.34))
    add_para(tf2, descriptor, 12.5, GREY, italic=True, first=True, space_after=0)
    # thin divider
    rect(s, Inches(0.45), Inches(0.98), Inches(12.45), Pt(1.6), fill=GOLD)
    # footer
    ftb, ftf = textbox(s, Inches(0.45), Inches(7.12), Inches(12.45), Inches(0.3))
    src = sources or "Source: company SEC filings (Form 10-K, CIK 0001397187)"
    add_para(ftf, src, 8.5, GREY, first=True, space_after=0)
    if page is not None:
        ptb, ptf = textbox(s, Inches(12.3), Inches(7.12), Inches(0.9), Inches(0.3))
        add_para(ptf, str(page), 9, GREY, align=PP_ALIGN.RIGHT, first=True, space_after=0)
    return s


def body_box(slide, l=Inches(0.5), t=Inches(1.12), w=Inches(12.35), h=Inches(5.9)):
    return textbox(slide, l, t, w, h)


PAGE = [0]
def pg():
    PAGE[0] += 1
    return PAGE[0]

# =====================================================================
# 1. TITLE
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, Inches(2.55), SW, Inches(0.06), fill=GOLD)
rect(s, 0, Inches(4.35), SW, Inches(0.06), fill=CARD)
tb, tf = textbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.4))
add_para(tf, "GLOBAL INVESTMENT SOCIETY", 20, GOLD, bold=True, first=True, space_after=2)
add_para(tf, "Investment Research Division", 15, WHITE, space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6))
add_para(tf, "lululemon athletica inc.", 40, WHITE, bold=True, first=True, space_after=0)
add_para(tf, "NASDAQ: LULU", 20, GOLD, bold=True, space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(4.6), Inches(11.6), Inches(2.2))
add_para(tf, "Recommendation:  OVERWEIGHT / LONG", 22, GOLD, bold=True, first=True, space_after=8)
add_para(tf, "Current price:  $100.00        Price target:  $140  (+40% upside)", 17, WHITE, space_after=6)
add_para(tf, "A net-cash, high-margin brand priced for terminal decline \u2014 we see a cyclical trough, not a broken business", 13.5, RGBColor(0xCF,0xD6,0xE4), italic=True, space_after=10)
add_para(tf, "September 2026        Prepared for the GIS IR selection process", 11, GREY, space_after=0)

# =====================================================================
# 2. TABLE OF CONTENTS
# =====================================================================
toc = [
    "1.  Investment thesis summary", "2.  Situation overview", "3.  Market narrative",
    "4.  Company overview", "5.  Business model & unit economics", "6.  Industry overview",
    "7.  Thesis I \u2014 Priced for terminal decline", "8.  Thesis II \u2014 International growth engine",
    "9.  Thesis III \u2014 Elite economics & capital return", "10. Risks & mitigants",
    "11. Catalyst timeline", "12. Financials \u2014 income statement",
    "13. Financials \u2014 balance sheet", "14. Financials \u2014 cash flow",
    "15. Capital structure & WACC", "16. Valuation summary (football field)",
    "17. DCF valuation", "18. Comparable companies", "19. Appendix \u2014 bull / bear scenarios",
]
s = slide_base("Table of Contents", "What this pitch will cover", page=pg())
tb, tf = textbox(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.8))
for i, item in enumerate(toc[:10]):
    add_para(tf, item, 14.5, NAVY, bold=True, first=(i == 0), space_after=8)
tb2, tf2 = textbox(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.8))
for j, item in enumerate(toc[10:]):
    add_para(tf2, item, 14.5, NAVY, bold=True, first=(j == 0), space_after=8)

# =====================================================================
# 3. INVESTMENT THESIS SUMMARY
# =====================================================================
s = slide_base("Investment Thesis Summary", "Overweight LULU with a $140 target on ~40% upside and an asymmetric payoff", page=pg())
# target box
box = rect(s, Inches(0.5), Inches(1.2), Inches(3.9), Inches(1.5), fill=LGREY)
btf = box.text_frame; btf.word_wrap = True
add_para(btf, "PRICE TARGET", 12, CARD, bold=True, first=True, space_after=2)
add_para(btf, "$140", 30, NAVY, bold=True, space_after=0)
add_para(btf, "+40% vs $100.00 today", 12, GREEN, bold=True, space_after=0)
# rating box
box2 = rect(s, Inches(4.6), Inches(1.2), Inches(3.9), Inches(1.5), fill=LGREY)
b2 = box2.text_frame; b2.word_wrap = True
add_para(b2, "RATING", 12, CARD, bold=True, first=True, space_after=2)
add_para(b2, "OVERWEIGHT", 24, NAVY, bold=True, space_after=0)
add_para(b2, "12-month horizon", 12, GREY, space_after=0)
# setup box
box3 = rect(s, Inches(8.7), Inches(1.2), Inches(4.15), Inches(1.5), fill=LGREY)
b3 = box3.text_frame; b3.word_wrap = True
add_para(b3, "WHY NOW", 12, CARD, bold=True, first=True, space_after=2)
add_para(b3, "Shares \u201355% off highs; ~18% single-day drop after Q2 FY2026 print (Sep 3, 2026) overshoots the fundamentals", 12.5, INK, space_after=0)
tb, tf = textbox(s, Inches(0.5), Inches(2.95), Inches(12.35), Inches(4.0))
add_para(tf, "Three reasons to be long", 15, CARD, bold=True, first=True, space_after=6)
add_para(tf, "1.  Priced for terminal decline \u2014 at ~3.5x EV/EBITDA and ~10x FY2026E EPS with a net-cash balance sheet, the stock embeds a permanent impairment that the business does not support", 13.5, INK, bold=False, space_after=7)
add_para(tf, "2.  International is a multi-year growth engine \u2014 China Mainland and Rest-of-World more than offset a maturing Americas and can return total revenue to mid-single-digit growth", 13.5, INK, space_after=7)
add_para(tf, "3.  Elite economics + accretive buybacks \u2014 FY2025 operating margin was 19.9% (last full year, not the trough); we model a 13.2% clean FY2026 run-rate (Q2 18.8% minus 560bps of tariff refunds), then add the $134.5M refund once (~14.5% reported). Diluted shares are down ~7% (128.0M \u2192 119.1M)", 13.5, INK, space_after=7)
add_para(tf, "Base-case DCF $124; bear $66 (\u221234%) vs bull $227 (+127%) \u2014 downside is protected by net cash and an ~8\u20139% FCF yield", 13, NAVY, bold=True, italic=True, space_after=0)

# =====================================================================
# 4. SITUATION OVERVIEW
# =====================================================================
s = slide_base("Situation Overview", "A high-quality compounder has de-rated to value multiples after a guidance reset", page=pg(),
               sources="Source: company 10-K (FY2025, ended Feb 1, 2026) and Q2 FY2026 release (Sep 3, 2026)")
tb, tf = body_box(s)
add_para(tf, "What happened", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "LULU shares fell from a 52-week high of ~$226 to ~$100, including an ~18% single-day decline following the Q2 FY2026 earnings release on September 3, 2026",
    "Management guided FY2026 to the first annual revenue decline in company history: net revenue of $10.35\u2013$10.50B (\u22125% to \u22127%) and diluted EPS of $9.48\u2013$9.73",
    "The Americas (the majority of sales) has decelerated to negative comparable sales, while tariffs and higher promotions pressure gross margin and SG&A deleverages on lower volume",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=5)
add_para(tf, "Where the business stands (FY2025, ended Feb 1, 2026)", 14.5, CARD, bold=True, space_after=5)
for t in [
    "Net revenue $11,102.6M (+4.9% y/y); gross profit $6,284.1M (56.6% margin); operating income $2,210.6M (19.9% margin)",
    "Net income $1,579.2M; diluted EPS $13.26; diluted shares 119.1M (down from 128.0M in FY2022)",
    "Balance sheet: $1,807.2M cash and no funded debt \u2014 a net-cash position that funds buybacks and international expansion",
    "Cash from operations $1,602.5M; capital expenditures $680.8M \u2192 ~$0.9B free cash flow even in a decelerating year",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=5)

# =====================================================================
# 5. MARKET NARRATIVE
# =====================================================================
s = slide_base("Market Narrative", "Sentiment has capitulated \u2014 the sell-side is cutting targets into the print", page=pg(),
               sources="Source: sell-side research notes (Sep 2026); company guidance")
tb, tf = body_box(s)
add_para(tf, "What the street is saying", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "Consensus rating has drifted to \u201cHold / Reduce\u201d with a rapidly falling target; recent cuts include Goldman Sachs and JPMorgan to $95, Morgan Stanley $93 (Underweight), Wells Fargo $95, and Bank of America to $122",
    "Bear case: Americas saturation and negative comps, tariff-driven gross-margin pressure, SG&A deleverage, and share loss to emerging brands (Alo, Vuori, On)",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=5)
add_para(tf, "Our differentiated view", 14.5, CARD, bold=True, space_after=5)
for t in [
    "The tape has moved from euphoria (30\u201340x earnings in 2023\u201324) to capitulation (~10x today) \u2014 expectations are now low enough that even a conservative recovery re-rates the stock",
    "At ~3.5x EV/EBITDA the market implicitly assumes revenue and margins decline in perpetuity; our base case only needs stabilization, not a return to peak growth",
    "Consensus is extrapolating one weak fiscal year; we underwrite the international runway and structural margins the street is discounting to zero",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=5)

# =====================================================================
# 6. COMPANY OVERVIEW
# =====================================================================
s = slide_base("Company Overview", "A vertically integrated, direct-to-consumer technical apparel brand", page=pg())
tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(7.2), Inches(5.9))
add_para(tf, "Business", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "Founded 1998 in Vancouver; designs technical athletic apparel and accessories for yoga, training, running and everyday wear",
    "Sells through company-operated stores, a direct e-commerce channel, and select wholesale \u2014 a predominantly DTC model that captures full retail economics and first-party data",
    "Reports revenue by geography: Americas (majority of sales), China Mainland (fastest-growing), and Rest of World",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)
add_para(tf, "Revenue trajectory (US$ M)", 14.5, CARD, bold=True, space_after=5)
for y in ["FY2022", "FY2023", "FY2024", "FY2025"]:
    add_para(tf, f"{y}:  {D.IS['revenue'][y]:,}   \u2192  op. margin {D.IS['operating_income'][y]/D.IS['revenue'][y]*100:.1f}%",
             13, INK, bullet=True, space_after=4)
# right: segment mix note box
box = rect(s, Inches(8.0), Inches(1.15), Inches(4.85), Inches(5.5), fill=LGREY)
btf = box.text_frame; btf.word_wrap = True
add_para(btf, "GROWTH ARCHITECTURE", 12.5, CARD, bold=True, first=True, space_after=6)
add_para(btf, "Americas", 13.5, NAVY, bold=True, space_after=1)
add_para(btf, "Mature, high-productivity base; now negative comps \u2014 the source of the market's concern", 12, INK, space_after=8)
add_para(btf, "China Mainland", 13.5, NAVY, bold=True, space_after=1)
add_para(btf, "Fastest-growing region (~20%+); large store and brand-awareness runway", 12, INK, space_after=8)
add_para(btf, "Rest of World", 13.5, NAVY, bold=True, space_after=1)
add_para(btf, "Early-stage in Europe and APAC; men's and international whitespace underpenetrated", 12, INK, space_after=0)

# =====================================================================
# 7. BUSINESS MODEL
# =====================================================================
s = slide_base("Business Model & Unit Economics", "Premium pricing plus DTC scale drives sector-leading margins", page=pg())
tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(7.2), Inches(5.9))
add_para(tf, "Unit economics", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "Gross margin ~57\u201359% \u2014 premium, full-price selling with historically limited promotion",
    "Operating margin ~20% (peaked ~24% in FY2024) \u2014 among the best in branded apparel",
    "Asset-efficient store fleet: leased footprint (right-of-use assets ~$1.6B) with high sales productivity per square foot",
    "Cash-generative: FY2025 CFO $1,602.5M on $680.8M capex \u2192 ~$0.9B free cash flow",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)
add_para(tf, "Competitive moats", 14.5, CARD, bold=True, space_after=5)
for t in [
    "Brand and community (ambassadors, in-store events) supporting pricing power",
    "Product innovation and proprietary fabrics driving repeat purchase",
    "First-party DTC data enabling assortment, pricing and inventory discipline",
    "Scale in sourcing and a growing men's and international opportunity",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)
box = rect(s, Inches(8.0), Inches(1.15), Inches(4.85), Inches(3.0), fill=NAVY)
btf = box.text_frame; btf.word_wrap = True
add_para(btf, "MARGIN PROFILE (FY2025)", 12.5, GOLD, bold=True, first=True, space_after=8)
add_para(btf, "Gross margin        56.6%", 15, WHITE, space_after=6)
add_para(btf, "Operating margin   19.9%", 15, WHITE, space_after=6)
add_para(btf, "Net margin            14.2%", 15, WHITE, space_after=6)
add_para(btf, "FCF (approx.)         ~$0.9B", 15, WHITE, space_after=0)

# =====================================================================
# 8. INDUSTRY OVERVIEW
# =====================================================================
s = slide_base("Industry Overview", "Structural tailwinds favor scaled brands in a consolidating category", page=pg(),
               sources="Source: company filings; industry estimates (analyst)")
tb, tf = body_box(s)
add_para(tf, "Category dynamics", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "Global activewear / athleisure is a large (~$400B) market growing at a mid-single-digit rate, supported by health, wellness and the ongoing casualization of apparel",
    "The category is fragmented but consolidating toward brands with scale, technical product and omnichannel reach; Nike and adidas are incumbents, with Vuori, Alo and On emerging",
    "Barriers to entry are rising: brand equity, technical fabric development, supply-chain scale and a productive retail footprint are difficult to replicate quickly",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)
add_para(tf, "Where LULU fits", 14.5, CARD, bold=True, space_after=5)
for t in [
    "A premium, category-defining brand in technical apparel with structurally higher margins than most peers",
    "Under-indexed internationally versus Nike/adidas \u2014 the primary multi-year growth lever",
    "Competition is real but LULU's scale, innovation cadence and community remain differentiated",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)

# =====================================================================
# 9-11. THESIS I / II / III
# =====================================================================
def thesis_slide(num, title, desc, headline, bullets, metric_title, metric_lines):
    s = slide_base(f"Investment Thesis {num}", desc, page=pg())
    tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(7.4), Inches(5.9))
    add_para(tf, headline, 15, CARD, bold=True, first=True, space_after=7)
    for t in bullets:
        add_para(tf, t, 13.3, INK, bullet=True, space_after=7)
    box = rect(s, Inches(8.15), Inches(1.15), Inches(4.7), Inches(len(metric_lines)*0.55 + 0.9), fill=LGREY)
    btf = box.text_frame; btf.word_wrap = True
    add_para(btf, metric_title, 12.5, CARD, bold=True, first=True, space_after=8)
    for ml in metric_lines:
        add_para(btf, ml, 13.5, NAVY, bold=True, space_after=6)
    return s

thesis_slide(
    "I", "Priced for terminal decline",
    "The valuation already discounts a permanently shrinking business",
    "The market is paying trough multiples for a brand that still earned 19.9% operating margin in FY2025",
    [
        "At ~$100 the stock trades at ~3.5x EV/EBITDA and ~10.4x FY2026E EPS \u2014 versus a 5-year history of ~20\u201330x earnings",
        "The balance sheet holds $1.8B of cash and no funded debt, so nearly the entire enterprise value is covered by the operating business at a very low multiple",
        "Our base-case DCF (WACC 10.5% from FRED 4.77% rf, terminal growth 2.25%, terminal EBIT margin 15.5%) is in the model \u2014 bear case still brackets the current price",
        "To justify $100 you must assume revenue and margins fall in perpetuity; that is inconsistent with international growth and the FY2026 tariff-refund tailwind",
    ],
    "VALUATION SNAPSHOT",
    ["EV / EBITDA:  3.5x", "FY2026E P/E:  10.4x", "FCF yield:  ~8\u20139%", "Net cash:  $1.8B", "DCF base:  $124"],
)

thesis_slide(
    "II", "International growth engine",
    "China and Rest-of-World offset a maturing Americas and restore growth",
    "Geographic mix shift is the bridge back to mid-single-digit revenue growth",
    [
        "China Mainland has been compounding at ~20%+ with a long runway in store count and brand awareness relative to the Americas base",
        "Rest of World (Europe, APAC) is early-stage and under-penetrated versus global peers \u2014 an incremental multi-year contributor",
        "As international scales, it dilutes the drag from negative Americas comps; our model returns total revenue to ~2.3% growth (Street 3Y forecast 2.26%) after the FY2026 reset",
        "Men's and the digital channel add further optionality that the market is not paying for today",
    ],
    "REVENUE PATH (model)",
    ["FY2025A:  $11,102.6M", "FY2026E:  $10,425M (\u22126.1%)", "FY2028E:  $10,942M", "FY2030E:  $11,451M"],
)

thesis_slide(
    "III", "Elite economics & capital return",
    "Best-in-class margins and buybacks compound value through the trough",
    "Even a down year throws off ~$1B of FCF that is being returned aggressively",
    [
        "Operating margin remains ~20% in FY2025 \u2014 far above most branded-apparel peers \u2014 with room to recover as tariffs and promotions normalize",
        "FY2025 cash from operations was $1,602.5M; after ~$681M capex the business generated roughly $0.9B of free cash flow",
        "The company repurchased $1.6B (FY2024) and $1.2B (FY2025) of stock; share count has fallen from 128.0M (FY2022) to 119.1M (FY2025)",
        "At ~$100 per share, every dollar of buyback retires far more shares than at prior highs \u2014 highly accretive to per-share value",
    ],
    "CAPITAL RETURN",
    ["FY2025 CFO:  $1,602.5M", "FY2025 capex:  $680.8M", "FY2025 buyback:  $1,178.3M", "Shares:  128.0M \u2192 119.1M"],
)

# =====================================================================
# 12. RISKS & MITIGANTS
# =====================================================================
s = slide_base("Risks & Mitigants", "The bear points are real but largely discounted at today's multiple", page=pg())
risks = [
    ("Brand fatigue / fashion risk in a core-heavy assortment",
     "Deep product pipeline, community engagement and international whitespace diversify demand"),
    ("Americas comparable sales stay negative",
     "Expectations are already low; ~10x earnings provides a valuation cushion if declines moderate"),
    ("Tariffs compress gross margin",
     "Pricing power, sourcing diversification, and $0.86/share of FY2026 tariff refunds offset part of the hit"),
    ("Competition from Alo, Vuori, On and incumbents",
     "Scale, fabric innovation, and under-penetrated men's / international segments defend share"),
    ("FX translation on a growing international mix",
     "Natural operational hedges and an active hedging program limit earnings volatility"),
]
tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(12.35), Inches(0.3))
add_para(tf, "Risk  \u2192  Mitigant", 14, CARD, bold=True, first=True, space_after=0)
top = 1.55
for rk, mg in risks:
    b = rect(s, Inches(0.5), Inches(top), Inches(6.0), Inches(0.95), fill=LGREY)
    bt = b.text_frame; bt.word_wrap = True; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(bt, rk, 12.5, CARD, bold=True, first=True, space_after=0)
    b2 = rect(s, Inches(6.7), Inches(top), Inches(6.15), Inches(0.95), fill=WHITE, line=GREY)
    bt2 = b2.text_frame; bt2.word_wrap = True; bt2.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(bt2, mg, 12.5, INK, first=True, space_after=0)
    top += 1.06

# =====================================================================
# 13. CATALYST TIMELINE
# =====================================================================
s = slide_base("Catalyst Timeline", "A sequence of events that closes the gap to intrinsic value", page=pg())
cats = [
    ("Q3 FY2026 (Dec 2026)", "Company guided revenue \u221210% to \u221211%; watch for stabilization signals and holiday traffic / China 11.11 read-through"),
    ("FY2026 year-end (early 2027)", "First full year lapping the reset; tariff refunds and cost actions support EPS versus lowered expectations"),
    ("FY2027", "Margin inflection as promotions/tariffs anniversary; continued accretive buybacks compound per-share value"),
    ("FY2027\u2013FY2028", "Total revenue returns to growth as international scales \u2014 the trigger for multiple re-rating toward peers"),
]
top = 1.5
for i, (when, what) in enumerate(cats):
    dot = rect(s, Inches(0.6), Inches(top+0.05), Inches(0.22), Inches(0.22), fill=CARD)
    b = rect(s, Inches(1.1), Inches(top), Inches(2.9), Inches(1.05), fill=NAVY)
    bt = b.text_frame; bt.word_wrap = True; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(bt, when, 13, GOLD, bold=True, first=True, space_after=0)
    b2 = rect(s, Inches(4.2), Inches(top), Inches(8.6), Inches(1.05), fill=LGREY)
    bt2 = b2.text_frame; bt2.word_wrap = True; bt2.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(bt2, what, 13, INK, first=True, space_after=0)
    top += 1.25

# =====================================================================
# Financial statement table helper
# =====================================================================
def stmt_table(slide, rows, headers, col0w=3.6, top=1.35, height=5.4, red_rows=(), bold_rows=()):
    ncol = len(headers)
    left = Inches(0.5)
    width = Inches(12.35)
    tbl_shape = slide.shapes.add_table(len(rows)+1, ncol, left, Inches(top), width, Inches(height))
    table = tbl_shape.table
    table.columns[0].width = Inches(col0w)
    restw = (12.35 - col0w) / (ncol - 1)
    for c in range(1, ncol):
        table.columns[c].width = Inches(restw)
    # style header
    for c, htxt in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = htxt
        _set_font(r, 11, WHITE, bold=True)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows, start=1):
        is_red = ri-1 in red_rows
        is_bold = ri-1 in bold_rows
        for c, val in enumerate(row):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LGREY
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = val
            color = CARD if (is_red and c > 0) else (NAVY if is_bold else INK)
            _set_font(r, 10.5, color, bold=is_bold)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
    return table

def m(v):
    return f"{v:,.0f}"

HY = ["FY2022", "FY2023", "FY2024", "FY2025"]
# projected values (from verified model)
PROJ_REV = {"FY2026E": 10425341, "FY2027E": 10529595, "FY2028E": 10845483, "FY2029E": 11279302, "FY2030E": 11730474}

# =====================================================================
# 14. INCOME STATEMENT
# =====================================================================
s = slide_base("Financials \u2014 Income Statement", "Historical results with our base-case forecast (US$ M)", page=pg(),
               sources="Source: company 10-K filings; projections per GIS operating model")
hdr = ["US$ M"] + HY + ["FY2026E", "FY2028E", "FY2030E"]
rows = [
    ["Net revenue"] + [m(D.IS['revenue'][y]/1000) for y in HY] + ["10,425", "10,845", "11,730"],
    ["Gross profit"] + [m(D.IS['gross_profit'][y]/1000) for y in HY] + ["6,025", "6,292", "6,642"],
    ["Operating income"] + [m(D.IS['operating_income'][y]/1000) for y in HY] + ["1,587", "1,853", "2,112"],
    ["Operating margin %"] + [f"{D.IS['operating_income'][y]/D.IS['revenue'][y]*100:.1f}%" for y in HY] + ["15.2%", "16.9%", "18.4%"],
    ["Net income"] + [m(D.IS['net_income'][y]/1000) for y in HY] + ["1,142", "1,322", "1,496"],
    ["Diluted EPS ($)"] + [f"{D.IS['diluted_eps'][y]:.2f}" for y in HY] + ["10.02", "12.58", "15.40"],
]
stmt_table(s, rows, hdr, col0w=3.2, bold_rows=(0, 2, 5))
tb, tf = textbox(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.35))
add_para(tf, "FY2026E reported OM 15.2% = clean ~13.9% (56.5% GM \u2212 42.5% SG&A) plus $134.5M IEEPA refund (~+1.3 ppt). DCF uses the Q2 clean 13.2% run-rate as its trough.", 11, GREY, italic=True, first=True, space_after=0)

# =====================================================================
# 15. BALANCE SHEET
# =====================================================================
s = slide_base("Financials \u2014 Balance Sheet", "A net-cash balance sheet underpins downside protection (US$ M)", page=pg(),
               sources="Source: company 10-K filings; projections per GIS operating model")
hdr = ["US$ M"] + HY + ["FY2026E", "FY2028E", "FY2030E"]
rows = [
    ["Cash & equivalents"] + [m(D.BS['cash'][y]/1000) for y in HY] + ["2,375", "3,778", "5,648"],
    ["Inventories"] + [m(D.BS['inventories'][y]/1000) for y in HY] + ["1,542", "1,535", "1,587"],
    ["Total assets"] + [m(D.BS['total_assets'][y]/1000) for y in HY] + ["8,946", "10,728", "12,888"],
    ["Total liabilities"] + [m(D.BS['total_liab'][y]/1000) for y in HY] + ["3,280", "3,432", "3,586"],
    ["Total equity"] + [m(D.BS['total_equity'][y]/1000) for y in HY] + ["5,666", "7,296", "9,303"],
    ["Funded debt"] + ["0", "0", "0", "0"] + ["0", "0", "0"],
]
stmt_table(s, rows, hdr, col0w=3.2, bold_rows=(2, 4, 5))
tb, tf = textbox(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.35))
add_para(tf, "No funded debt and a growing cash balance fund buybacks and international expansion \u2014 the core of our margin of safety", 11, GREY, italic=True, first=True, space_after=0)

# =====================================================================
# 16. CASH FLOW
# =====================================================================
s = slide_base("Financials \u2014 Cash Flow", "Durable free cash flow through the trough (US$ M)", page=pg(),
               sources="Source: company 10-K filings; projections per GIS operating model")
hdr = ["US$ M"] + HY + ["FY2026E", "FY2028E", "FY2030E"]
rows = [
    ["Cash from operations"] + [m(D.CF['cfo'][y]/1000) for y in HY] + ["1,797", "1,894", "2,068"],
    ["Capital expenditures"] + [f"({m(D.CF['capex'][y]/1000)})" for y in HY] + ["(730)", "(602)", "(573)"],
    ["Free cash flow"] + [m((D.CF['cfo'][y]-D.CF['capex'][y])/1000) for y in HY] + ["1,067", "1,292", "1,495"],
    ["Share repurchases"] + [f"({m(D.CF['buybacks'][y]/1000)})" for y in HY] + ["(500)", "(500)", "(500)"],
    ["D&A"] + [m(D.CF['d_and_a'][y]/1000) for y in HY] + ["480", "492", "515"],
]
stmt_table(s, rows, hdr, col0w=3.2, bold_rows=(2,))
tb, tf = textbox(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.35))
add_para(tf, "Free cash flow stays around $1B+ even in the FY2026 reset year, comfortably funding continued repurchases", 11, GREY, italic=True, first=True, space_after=0)

# =====================================================================
# 17. CAPITAL STRUCTURE & WACC
# =====================================================================
s = slide_base("Capital Structure & WACC", "A ~100% equity, net-cash structure drives a ~10.5% discount rate", page=pg())
tb, tf = textbox(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.6))
add_para(tf, "Capital structure", 14.5, CARD, bold=True, first=True, space_after=5)
for t in [
    "No funded debt; an undrawn revolving credit facility provides liquidity",
    "$1,807.2M cash and equivalents \u2014 a net-cash position",
    "Enterprise value \u2248 equity value less cash; ~$9.3B EV at ~$100 per share",
    "We use Damodaran\u2019s unlevered Retail (Special Lines) beta of 0.95 because LULU has no debt, so \u03b2u = \u03b2e. Company 5Y is 0.86 and is not used.",
    "Capital returned via buybacks (no dividend); FY2025 repurchases $1,178.3M",
]:
    add_para(tf, t, 13, INK, bullet=True, space_after=6)
box = rect(s, Inches(7.0), Inches(1.2), Inches(5.85), Inches(4.6), fill=LGREY)
btf = box.text_frame; btf.word_wrap = True
add_para(btf, "WACC BUILD (CAPM)", 13, CARD, bold=True, first=True, space_after=8)
for t, v in [
    ("Risk-free rate (10-yr UST)", "4.8%"),
    ("Equity risk premium", "6.0%"),
    ("Beta (unlev. retail; no debt)", "0.95"),
    ("Cost of equity", "10.5%"),
    ("Debt weight", "0%"),
    ("Equity weight", "100%"),
]:
    p = btf.add_paragraph(); p.space_after = Pt(6)
    r = p.add_run(); r.text = f"{t}"; _set_font(r, 13, INK, bold=(t in ("Cost of equity",)))
    r2 = p.add_run(); r2.text = f"      {v}"; _set_font(r2, 13, NAVY, bold=True)
p = btf.add_paragraph(); p.space_before = Pt(4)
r = p.add_run(); r.text = "WACC = 10.5%"; _set_font(r, 16, GREEN, bold=True)

# =====================================================================
# 18. VALUATION SUMMARY (FOOTBALL FIELD)
# =====================================================================
s = slide_base("Valuation Summary", "Multiple methods converge above the current price \u2014 target $140", page=pg(),
               sources="Source: GIS DCF and comps models; multiples are analyst ranges")
# football field: horizontal floating bars
methods = [
    ("P / E (10\u201318x FY2026E)", 96, 173),
    ("EV / EBITDA (4.5\u20137.5x)", 126, 199),
    ("DCF (bear \u2013 bull)", 91, 239),
    ("52-week range", 100, 226),
]
chart_l, chart_r = 3.2, 12.6
vmin, vmax = 60, 250
def xpos(v):
    return chart_l + (v - vmin) / (vmax - vmin) * (chart_r - chart_l)
top = 1.7
for name, lo, hi in methods:
    tb, tf = textbox(s, Inches(0.5), Inches(top-0.02), Inches(2.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, name, 12, NAVY, bold=True, first=True, space_after=0)
    bar = rect(s, Inches(xpos(lo)), Inches(top), Inches(xpos(hi)-xpos(lo)), Inches(0.45), fill=GOLD)
    bt = bar.text_frame; bt.word_wrap = False; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    bt.margin_left = Pt(4); bt.margin_right = Pt(4)
    p = bt.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"${lo}"; _set_font(r, 10.5, NAVY, bold=True)
    tb2, tf2 = textbox(s, Inches(xpos(hi)+0.02), Inches(top), Inches(0.9), Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf2, f"${hi}", 10.5, NAVY, bold=True, first=True, space_after=0)
    top += 0.75
# current price line and target line
cp_x = xpos(100); pt_x = xpos(140)
ln = rect(s, Inches(cp_x), Inches(1.55), Pt(2), Inches(3.4), fill=INK)
ln2 = rect(s, Inches(pt_x), Inches(1.55), Pt(2), Inches(3.4), fill=CARD)
tb, tf = textbox(s, Inches(cp_x-0.7), Inches(4.95), Inches(1.6), Inches(0.3))
add_para(tf, "Current $100", 10.5, INK, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
tb, tf = textbox(s, Inches(pt_x-0.7), Inches(5.2), Inches(1.6), Inches(0.3))
add_para(tf, "Target $140", 10.5, CARD, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
tb, tf = textbox(s, Inches(0.5), Inches(5.7), Inches(12.35), Inches(1.2))
add_para(tf, "We set a 12-month target of $140 \u2014 above the conservative Gordon DCF of $124, for a partial re-rating toward historical multiples \u2014 implying ~40% upside", 13, NAVY, bold=True, first=True, space_after=5)
add_para(tf, "The P/E low end (~$96) sits near today's price, showing how little recovery is required for the stock to work", 12.5, INK, italic=True, space_after=0)

# =====================================================================
# 19. DCF VALUATION
# =====================================================================
s = slide_base("DCF Valuation", "Base-case unlevered DCF yields ~$124 per share", page=pg(),
               sources="Source: GIS DCF model (from scratch); FY2025 cash & share count per 10-K")
# assumptions table
tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(6.1), Inches(0.4))
add_para(tf, "Base-case assumptions", 14, CARD, bold=True, first=True, space_after=0)
arows = [
    ["Assumption", "Value"],
    ["Revenue growth (FY26 \u2192 FY30)", "\u22126.1% \u2192 +2.3%"],
    ["EBIT margin (FY26 \u2192 FY30)", "13.2% clean \u2192 15.5% (+$134.5M FY26)"],
    ["Tax rate", "30%"],
    ["Capex % of revenue", "5.5% blend"],
    ["WACC", "10.5%"],
    ["Terminal growth", "2.25%"],
]
t = s.shapes.add_table(len(arows), 2, Inches(0.5), Inches(1.6), Inches(6.0), Inches(3.0)).table
t.columns[0].width = Inches(4.0); t.columns[1].width = Inches(2.0)
for ri, row in enumerate(arows):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci); cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (WHITE if ri % 2 else LGREY)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = val
        _set_font(r, 11.5, WHITE if ri == 0 else (CARD if ci == 1 and ri > 0 else INK), bold=(ri == 0))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
# output box
box = rect(s, Inches(7.0), Inches(1.6), Inches(5.85), Inches(3.0), fill=LGREY)
btf = box.text_frame; btf.word_wrap = True
add_para(btf, "VALUATION OUTPUT (US$ M)", 12.5, CARD, bold=True, first=True, space_after=6)
for t2, v in [
    ("PV of explicit FCF (FY26\u2013FY30)", "3,732"),
    ("PV of terminal value", "8,316"),
    ("Enterprise value", "12,048"),
    ("Plus: cash", "1,807"),
    ("Equity value", "13,856"),
    ("\u00f7 Diluted shares (M)", "111.4"),
]:
    p = btf.add_paragraph(); p.space_after = Pt(5)
    r = p.add_run(); r.text = t2; _set_font(r, 12.5, INK, bold=("Enterprise" in t2 or "Equity" in t2))
    r2 = p.add_run(); r2.text = f"      {v}"; _set_font(r2, 12.5, NAVY, bold=True)
p = btf.add_paragraph(); p.space_before = Pt(6)
r = p.add_run(); r.text = "Implied value:  $124 / share"; _set_font(r, 16, GREEN, bold=True)
# sensitivity mini
tb, tf = textbox(s, Inches(0.5), Inches(4.85), Inches(12.35), Inches(2.0))
add_para(tf, "Sensitivity \u2014 implied share price (WACC vs terminal growth)", 13, CARD, bold=True, first=True, space_after=4)
sens = [
    ["WACC \\ g", "1.5%", "2.0%", "2.25%", "2.5%", "3.0%"],
    ["9.5%", "$131", "$136", "$140", "$143", "$151"],
    ["10.5%", "$118", "$122", "$124", "$127", "$132"],
    ["11.5%", "$107", "$111", "$113", "$114", "$119"],
]
st = s.shapes.add_table(len(sens), 6, Inches(0.5), Inches(5.35), Inches(7.6), Inches(1.5)).table
for ri, row in enumerate(sens):
    for ci, val in enumerate(row):
        cell = st.cell(ri, ci); cell.fill.solid()
        hot = (ri == 2 and ci == 3)
        cell.fill.fore_color.rgb = NAVY if (ri == 0 or ci == 0) else (GOLD if hot else (WHITE if ri % 2 else LGREY))
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        _set_font(r, 10.5, WHITE if (ri == 0 or ci == 0) else (NAVY if hot else INK), bold=(ri == 0 or ci == 0 or hot))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# =====================================================================
# 20. COMPS
# =====================================================================
s = slide_base("Comparable Companies", "LULU screens cheap vs public peers; Gymshark is a private growth print, not the TV", page=pg(),
               sources="Source: StockAnalysis public prints; Gymshark 23.5x = Guardian £1.25bn / SGB FY25 EBITDA £53.3m")
hdr = ["Company", "EV/EBITDA", "P/E (fwd)", "Rev growth", "Op margin"]
rows = [
    ["lululemon (LULU)", "3.5x", "10.4x", "\u22126% (FY26E)", "~20%"],
    ["Nike (NKE)", "12.0x", "22.4x", "low-single", "~11%"],
    ["Deckers (DECK)", "8.0x", "11.2x", "mid-teens", "~22%"],
    ["On Holding (ONON)", "14.0x", "15.2x", "20%+", "~10%"],
    ["adidas (ADS)", "9.3x", "14.2x", "mid-single", "~9%"],
    ["V.F. Corp (VFC)", "10.7x", "11.7x", "flat/decl.", "~8%"],
    ["Gymshark (private)", "23.5x", "n.m.", "low-single", "~8% EBITDA"],
]
stmt_table(s, rows, hdr, col0w=3.4, top=1.25, height=3.85, bold_rows=(0,))
tb, tf = textbox(s, Inches(0.5), Inches(5.2), Inches(12.35), Inches(1.75))
add_para(tf, "Read-through \u2014 we do not average these for terminal value", 13.5, CARD, bold=True, first=True, space_after=3)
add_para(tf, "Public 5-name mean is 10.8x; with Gymshark 12.9x. Both overstate a 2.25% g / 15.5% OM FY2030 exit", 12.5, INK, bullet=True, space_after=3)
add_para(tf, "Selected TV is 8.0x = Deckers, ~1 turn above Gordon ~7x and a haircut from the mature public mean (DECK / ADS / VFC) of 9.3x", 12.5, INK, bullet=True, space_after=3)
add_para(tf, "Gymshark 23.5x = £1,250m / £53.3m is a 2020 GA print on FY25 EBITDA \u2014 a growth-stage private reference, not a mature public exit", 12.5, INK, bullet=True, space_after=0)

# =====================================================================
# 21. APPENDIX: SCENARIOS
# =====================================================================
s = slide_base("Appendix \u2014 Bull / Bear Scenarios", "Asymmetric payoff: limited downside, substantial upside", page=pg(),
               sources="Source: GIS DCF model (scenario tab)")
cols = [
    ("BEAR", "$66", "\u221234%", CARD, [
        "FY2026 revenue \u22129%; growth stays negative (\u22121% avg FY28\u201330)",
        "Terminal EBIT margin 12.0%",
        "WACC 11.5%; terminal growth 1.5%",
        "Americas decline persists; share loss continues",
    ]),
    ("BASE", "$124", "+24%", GREEN, [
        "FY2026 revenue \u22126.1%; then +2.3% (Street 3Y forecast)",
        "Clean EBIT margin 13.2% \u2192 15.5%; +$134.5M refund in FY26 only",
        "WACC 10.5%; terminal growth 2.25%",
        "International offsets a stabilizing Americas",
    ]),
    ("BULL", "$227", "+127%", NAVY, [
        "FY2026 revenue \u22124%; +6% avg FY27\u201330",
        "Terminal EBIT margin 19.0%",
        "WACC 9.5%; terminal growth 3.0%",
        "Margin recovery toward peak; brand re-accelerates",
    ]),
]
x = 0.5
for name, px, up, color, bullets in cols:
    head = rect(s, Inches(x), Inches(1.35), Inches(4.05), Inches(0.95), fill=color)
    ht = head.text_frame; ht.word_wrap = True; ht.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_para(ht, f"{name}   {px}", 20, WHITE, bold=True, first=True, space_after=0, align=PP_ALIGN.CENTER)
    add_para(ht, f"{up} vs $100", 12.5, WHITE, align=PP_ALIGN.CENTER, space_after=0)
    b = rect(s, Inches(x), Inches(2.4), Inches(4.05), Inches(3.6), fill=LGREY)
    bt = b.text_frame; bt.word_wrap = True
    for i, blt in enumerate(bullets):
        add_para(bt, blt, 12.5, INK, bullet=True, first=(i == 0), space_after=8)
    x += 4.25
tb, tf = textbox(s, Inches(0.5), Inches(6.2), Inches(12.35), Inches(0.8))
add_para(tf, "Probability-weighted value (25% / 50% / 25%) \u2248 $155 \u2014 the risk/reward skews decisively to the upside", 13.5, NAVY, bold=True, italic=True, first=True, space_after=0)

# =====================================================================
# 22. DISCLAIMER / SOURCES
# =====================================================================
s = slide_base("Sources & Disclaimer", "Data provenance and standard research disclaimer", page=pg())
tb, tf = body_box(s)
add_para(tf, "Sources", 14, CARD, bold=True, first=True, space_after=5)
for t in [
    "Historical financials: lululemon athletica inc. Forms 10-K via SEC EDGAR (CIK 0001397187); FY2025 fiscal year ended February 1, 2026",
    "Q2 FY2026 results and FY2026 guidance: company earnings release dated September 3, 2026",
    "Market data (price, shares, beta): public market sources as of early September 2026",
    "Projections, DCF and comparable-company analysis: GIS Investment Research models, built from scratch for this assignment",
    "Gymshark private print: Guardian (2020 GA £1.25bn valuation) and SGB (FY25 EBITDA £53.3m) \u2014 23.5x is a reference, not the FY30 exit",
]:
    add_para(tf, t, 12.5, INK, bullet=True, space_after=5)
add_para(tf, "Disclaimer", 14, CARD, bold=True, space_after=5)
add_para(tf, "This presentation is prepared for educational purposes as part of the Global Investment Society selection process and does not constitute investment advice or a recommendation to buy or sell any security", 12, GREY, italic=True, space_after=0)

prs.save(OUT)
print("Saved", os.path.abspath(OUT), "with", len(prs.slides._sldIdLst), "slides")
